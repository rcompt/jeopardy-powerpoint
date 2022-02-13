import argparse
import os
import random

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

class JeopardyBuilder:

    def __init__(self, template):
        self.template = template

        self.filename = "jeopardy-slides.pptx"
        
    def process_gameboard(self, slide, categories):
        for j, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            if shape.text.startswith('$'):
                continue
            shape.text = categories[shape.text]
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def process_DDs(self, board, DD_slides):
        DD_cards = random.sample([c for c in board.shapes 
                                if c.has_text_frame and c.text.startswith('$')], 
                                len(DD_slides))
        for i, card in enumerate(DD_cards):
            DD_slide_id, DD_slide = DD_slides[i]
            for shape in DD_slide.shapes:
                if not shape.has_text_frame:
                    # The DD image, and the sound icon
                    continue
                if 'Rectangle' not in shape.name:
                    # The Oval covering the sound icon
                    continue
                shape.click_action.hyperlink.address = card.click_action.hyperlink.address
                shape.click_action.target_slide = card.click_action.target_slide
                card.click_action.hyperlink.address = 'slide{}.xml'.format(DD_slide_id+1)
                card.click_action.target_slide = DD_slide
                break


    def build(self, categories, questions):

        assert os.path.exists(self.template), 'provided `jeopardy_ppt` does not exist.'

        prs = Presentation(self.template)

        game_boards = []
        daily_doubles = []
        for i, slide in enumerate(prs.slides):
            if i == 0:
                # Title slide
                continue
            elif i in (1, 2):
                # Game boards for Rounds 1 and 2
                self.process_gameboard(slide, categories)
                game_boards.append(slide)
            elif i == 3:
                # Final Jeopardy
                for j, shape in enumerate(slide.shapes):
                    if shape.click_action.hyperlink.address is not None:
                        continue
                    if not shape.has_text_frame:
                        continue
                    if j == 0:
                        # Always the title
                        pass
                    else:
                        shape.text = questions[shape.text]
                        shape.text_frame.word_wrap = True

            elif i in (4, 5, 6):
                # Daily doubles 
                daily_doubles.append((i, slide))
            else:
                if len(slide.shapes) != 4:
                    raise RuntimeError("Don't edit the fucking template! Need only 4 shapes per "
                                    "question slide.")
                num_text = 0
                for j, shape in enumerate(slide.shapes):
                    if shape.click_action.hyperlink.address is not None:
                        continue
                    if not shape.has_text_frame:
                        continue
                    num_text += 1
                    if j == 0:
                        # Always the title
                        pass
                    else:
                        # Always where the question goes
                        shape.text = questions[shape.text]
                        shape.text_frame.word_wrap = True
                if num_text > 2:
                    raise RuntimeError("Don't edit the fucking template! Only 2 text frames allowed "
                                    "per question card.")
                # ONETIME
                # slide.notes_slide.notes_text_frame.text = ""


        self.process_DDs(game_boards[0], daily_doubles[:1])
        self.process_DDs(game_boards[1], daily_doubles[1:])

        prs.save(self.filename)

        return self.filename

if __name__ == '__main__':
    builder = JeopardyBuilder(
        template = "../templates/jeopardy_template.pptx"
    )