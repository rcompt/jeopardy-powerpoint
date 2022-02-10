import argparse
import os
import pandas as pd
import random
import re
import shutil
import yaml

from pptx import Presentation
from pptx.enum.text import PP_ALIGN


def load_questions(infile):
    with open(infile) as iff:
        questions = yaml.load(iff, Loader=yaml.FullLoader)
    if len(questions) != 5:
        raise RuntimeError('Need 5 categories per round. For a blank category that will be filled '
                           'with music or pictures, Create a category with the expected name and '
                           'use keyword `BLANK` for all questions.')
    if not all([len(v) == 5 for v in questions.values()]):
        raise RuntimeError('All categories need to have 5 questions. Use `BLANK` to leave a '
                           'question blank on purpose.')
    questions = pd.DataFrame(questions)
    questions = questions[random.sample(list(questions.columns), 5)]
    return questions

def process_gameboard(slide, questions):
    i = 0
    for j, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        if shape.text.startswith('$'):
            continue
        shape.text = questions.columns[i]
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        i += 1

def process_DDs(board, DD_slides):
    DD_cards = random.sample([c for c in board.shapes 
                              if c.has_text_frame and c.text.startswith('$')], 
                             len(DD_slides))
    for i, card in enumerate(DD_cards):
        DD_slide_id, DD_slide = DD_slides[i]
        for shape in DD_slide.shapes:
            if not shape.has_text_frame:
                # The DD image, and the sound icon
                continue
            if 'Rectangle' not in shape.name:
                # The Oval covering the sound icon
                continue
            shape.click_action.hyperlink.address = card.click_action.hyperlink.address
            shape.click_action.target_slide = card.click_action.target_slide
            card.click_action.hyperlink.address = 'slide{}.xml'.format(DD_slide_id+1)
            card.click_action.target_slide = DD_slide
            break


def main():
    
    parser = argparse.ArgumentParser()
    parser.add_argument('--r1_questions', '-R1', help='Path to questions for round 1', required=True)
    parser.add_argument('--r2_questions', '-R2', help='Path to questions for round 2', required=True)
    parser.add_argument('--outfile', '-O', help='Path to the outfile', required=False, default=None)
    parser.add_argument('jeopardy_ppt', help='Path to the template')
    params = parser.parse_args()

    assert os.path.exists(params.jeopardy_ppt), 'provided `jeopardy_ppt` does not exist.'
    assert os.path.exists(params.r1_questions), 'provided `r1_questions` does not exist.'
    assert os.path.exists(params.r2_questions), 'provided `r2_questions` does not exist.'

    if params.outfile is not None:
        assert not os.path.exists(params.outfile), 'provided `outfile` exists.'
        outfile = params.outfile
    else:
        from datetime import datetime
        now = datetime.strftime(datetime.now(), '%m_%d_%Y_%H_%M_%S')
        outfile = '{}_{}{}'.format(os.path.splitext(params.jeopardy_ppt)[0], 
                                                    now,
                                                    os.path.splitext(params.jeopardy_ppt)[1])

    prs = Presentation(params.jeopardy_ppt)

    questions = {
        1: load_questions(params.r1_questions),
        2: load_questions(params.r2_questions)
    }

    board_re = re.compile(r'^R(?P<Rnum>([12])):'  # Round number
                          r'(?P<Cid>([0-4])),'   # Category Index
                          r'(?P<Qid>([0-4]))$')  # Question Index
    game_boards = []
    daily_doubles = []
    for i, slide in enumerate(prs.slides):
        if i == 0:
            # Title slide
            continue
        elif i in (1, 2):
            # Game boards for Rounds 1 and 2
            process_gameboard(slide, questions[i])
            game_boards.append(slide)
        elif i == 3:
            # Final Jeopardy. Do this manually
            continue
        elif i in (4, 5, 6):
            # Daily doubles 
            daily_doubles.append((i, slide))
        else:
            if len(slide.shapes) != 4:
                raise RuntimeError("Don't edit the fucking template! Need only 4 shapes per "
                                   "question slide.")
            num_text = 0
            for j, shape in enumerate(slide.shapes):
                if shape.click_action.hyperlink.address is not None:
                    continue
                if not shape.has_text_frame:
                    continue
                num_text += 1
                if j == 0:
                    # Always the title
                    q = board_re.search(shape.text)
                    assert(q is not None)
                    question = questions[int(q['Rnum'])].iloc[int(q['Cid']), int(q['Qid'])]
                else:
                    # Always where the question goes
                    if question == 'BLANK':
                        shape.text = "THIS QUESTION INTENTIONALLY LEFT BLANK"
                    else:
                        shape.text = question
                    shape.text_frame.word_wrap = True
            if num_text > 2:
                raise RuntimeError("Don't edit the fucking template! Only 2 text frames allowed "
                                   "per question card.")
            # ONETIME
            # slide.notes_slide.notes_text_frame.text = ""


    process_DDs(game_boards[0], daily_doubles[:1])
    process_DDs(game_boards[1], daily_doubles[1:])

    prs.save(outfile)

if __name__ == '__main__':
    main()